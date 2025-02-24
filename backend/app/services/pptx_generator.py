import os
import re
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from app.services.faiss_vector_db import clear_vector_db


def wrap_text_to_lines(text, line_length=80):
    """
    Naive word-wrap: splits 'text' into lines, each up to 'line_length' characters.
    This helps approximate how many lines we'll need for a given text box width.
    """
    if not text:
        return [""]

    words = text.split()
    lines = []
    current_line = []
    current_len = 0

    for word in words:
        extra_space = 1 if current_line else 0
        if current_len + len(word) + extra_space <= line_length:
            current_line.append(word)
            current_len += len(word) + extra_space
        else:
            lines.append(" ".join(current_line))
            current_line = [word]
            current_len = len(word)
    if current_line:
        lines.append(" ".join(current_line))

    return lines


def force_min40_max50_break(title_text):
    """
    1) If total words < 40 => one line
    2) If 40 >= total words <= 50 => two lines, first line has 40 words
    3) If total words > 50 => two lines, first line has 50 words
    """
    words = title_text.split()
    n = len(words)

    if n == 0:
        return ""

    if n < 40:
        # Keep all words in one line
        return " ".join(words)

    elif 40 >= n <= 50:
        line1 = " ".join(words[:40])
        line2 = " ".join(words[40:])
        return (line1 + "\n" + line2).strip()

    else:  # n > 50
        line1 = " ".join(words[:50])
        line2 = " ".join(words[50:])
        return (line1 + "\n" + line2).strip()


def add_auto_height_textbox(
    slide,
    left_in,
    top_in,
    width_in,
    max_height_in,
    text,
    font_size=24,
    align_center=True,
    bold=False,
    template_index=None,
):
    """
    Adds a text box with a fixed width and auto-calculated height based on the text length.
    """
    left = Inches(left_in)
    top = Inches(top_in)
    width = Inches(width_in)
    max_height = Inches(max_height_in)

    # If we're NOT center-aligning, use a large wrap width so PPT does normal paragraph wrapping
    if align_center:
        wrap_width = 80
    else:
        wrap_width = 999

    paragraphs = text.split("\n")  # handle forced newlines
    lines_for_render = []
    for paragraph_text in paragraphs:
        wrapped = wrap_text_to_lines(paragraph_text, line_length=wrap_width)
        if not wrapped:
            lines_for_render.append("")
        else:
            lines_for_render.extend(wrapped)

    lines = lines_for_render

    # Calculate shape height needed
    line_spacing_pt = font_size * 1.2
    line_height_emu = int(line_spacing_pt * 12700)  # 1 pt = 12700 EMU
    total_lines = len(lines)
    needed_height_emu = total_lines * line_height_emu

    # NO TRUNCATION: we do NOT cut off lines if shape is taller than max_height
    shape_height = needed_height_emu

    textbox = slide.shapes.add_textbox(left, top, width, shape_height)
    tf = textbox.text_frame
    tf.clear()

    # Fill the text box with paragraphs
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        if align_center:
            p.alignment = PP_ALIGN.CENTER
        else:
            p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold

            if template_index is not None:
                if template_index == 0:
                    run.font.name = "Arial"
                    run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
                elif template_index == 1:
                    run.font.name = "Times New Roman"
                    run.font.color.rgb = RGBColor(0x80, 0x00, 0x80)
                elif template_index == 2:
                    run.font.name = "Verdana"
                    run.font.color.rgb = RGBColor(0x00, 0x00, 0x8B)
                elif template_index == 3:
                    run.font.name = "Courier New"
                    run.font.color.rgb = RGBColor(0x8B, 0x45, 0x13)
                elif template_index == 4:
                    run.font.name = "Century Gothic"
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                elif template_index == 5:
                    run.font.name = "Palatino Linotype"
                    run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    tf.auto_size = None  # Prevent auto-resizing to a single line
    tf.word_wrap = True  # Let PPT handle final wrapping

    return textbox


def place_first_slide(
    slide, slide_width_in, slide_height_in, title_text, subtitle_text, template_index
):
    """
    Creates two text boxes on the first slide:
      1. Title at the top center (bold, large font), forced to two lines.
      2. Subtitle in the vertical middle (centered).
    """
    forced_title = force_min40_max50_break(title_text)
    width_in = 9.5
    title_left = (slide_width_in - width_in) / 2
    title_top = 1.5
    add_auto_height_textbox(
        slide=slide,
        left_in=title_left,
        top_in=title_top,
        width_in=width_in,
        max_height_in=2.0,
        text=forced_title,
        font_size=34,
        align_center=True,
        bold=True,
        template_index=template_index,
    )

    subtitle_left = title_left
    subtitle_top = slide_height_in / 2
    add_auto_height_textbox(
        slide=slide,
        left_in=subtitle_left,
        top_in=subtitle_top,
        width_in=width_in,
        max_height_in=3.0,
        text=subtitle_text,
        font_size=26,
        align_center=True,
        bold=False,
        template_index=template_index,
    )


def place_slide_content(
    slide, slide_width_in, title_text, content_text, template_index
):
    """
    For slides 2, 3, etc., places:
      - A larger title box near the top (bold)
      - A content box below it
    """
    width_in = 9.5
    left_in = (slide_width_in - width_in) / 2

    add_auto_height_textbox(
        slide=slide,
        left_in=left_in,
        top_in=1.0,
        width_in=width_in,
        max_height_in=1.5,
        text=title_text,
        font_size=32,
        align_center=True,
        bold=True,
        template_index=template_index,
    )

    add_auto_height_textbox(
        slide=slide,
        left_in=left_in,
        top_in=1.9,
        width_in=width_in,
        max_height_in=4.0,
        text=content_text,
        font_size=24,
        align_center=False,
        bold=False,
        template_index=template_index,
    )


def create_presentation(data, index):
    """
    Takes a data whose structure is as follows:
    {
        "title": '',
        "subtitle": '',
        "slides": [
            {
                "id": 1,
                "title": "",
                "content": "",
            }
            ...
        ]
    }

    Creates a presentation using 'test.pptx' as a template.
    - First slide: Title forcibly broken into lines, no truncation
    - Subsequent slides: Each has a bold title at top, content below
    """
    index_file = [
        "template0.pptx",
        "template1.pptx",
        "template2.pptx",
        "template3.pptx",
        "template4.pptx",
        "template5.pptx",
    ]
    # index = index
    try:
        # Ensure the index is within the valid range (0 to 5)
        if index < 0 or index > 5:
            raise IndexError(
                f"Index {index} is out of range. It must be between 0 and 5."
            )

        # Construct the file path
        print(f"Index Value: {index}")
        template_path = os.path.join(os.getcwd(), "app", "templates", index_file[index])
        print(f"File path: {template_path}")

        # Check if file exists
        if not os.path.exists(template_path):
            raise FileNotFoundError(f"Template file not found: {template_path}")

        # Load the presentation from temporary file
        prs = Presentation(template_path)
        print(f'Actual file path: {template_path}')

        if len(data.get("slides")) >= 15:
            print(
                "Number of slides in response is greater than the actual number of slides."
            )
            return None

        slide_width = prs.slide_width if prs.slide_width else 9144000
        slide_height = prs.slide_height if prs.slide_height else 6858000
        slide_width_in = slide_width / 914400.0
        slide_height_in = slide_height / 914400.0

        if len(prs.slides) == 0:
            blank_layout = prs.slide_layouts[-1]
            prs.slides.add_slide(blank_layout)

        first_slide = prs.slides[0]
        place_first_slide(
            first_slide,
            slide_width_in,
            slide_height_in,
            data.get("title", ""),
            data.get("subtitle", ""),
            index,
        )

        for i, slide_data in enumerate(data.get("slides", []), start=1):
            if i < len(prs.slides):
                slide = prs.slides[i]
            else:
                blank_layout = prs.slide_layouts[-1]
                slide = prs.slides.add_slide(blank_layout)

            place_slide_content(
                slide,
                slide_width_in,
                slide_data.get("title", ""),
                slide_data.get("content", ""),
                index,
            )

        # Remove any extra slides if needed
        for i in range(len(prs.slides) - 1, len(data.get("slides")), -1):
            rId = prs.slides._sldIdLst[i].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[i]

        pres_title = data.get("title", "Untitled Presentation")
        sanitized_title = re.sub(r'[<>:"\\|?*]', "", pres_title).strip()
        sanitized_title = sanitized_title.replace("/", "-").replace("\\", "-")
        filename = f"app/presentations/ppt/{sanitized_title}.pptx"
        os.makedirs(os.path.dirname(filename), exist_ok=True)
        prs.save(filename)

    except Exception as e:
        print(f"Error creating presentation: {e}")
        return None

    clear_vector_db()

    print(f"Presentation saved at: {filename}")
    return filename
