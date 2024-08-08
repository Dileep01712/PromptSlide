from pptx import Presentation  # type: ignore
from pptx.dml.color import RGBColor  # type: ignore
from pptx.util import Inches, Pt  # type: ignore


# Function to set background color
def set_background_color(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*color)


# Function to set font color and size for all text in the slide
def set_text_color_and_size(slide, color, title_size, para_size):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(*color)
                    if shape == slide.shapes.title:
                        run.font.size = Pt(title_size)
                    else:
                        run.font.size = Pt(para_size)


def create_presentation(filename, data):
    """
    Takes a data whose structure is as follows:
    {
        "title": '',
        "subtitle": '',
        "slides": [
            {
                "id": 1,
                "title": "",
                "content": "",
            }
            ....
        ]
    }
    Creates a presentation in the current directory based on the data.
    """
    # Create a Presentation object
    presentation = Presentation()

    # Add Title Slide
    title_slide_layout = presentation.slide_layouts[
        0
    ]  # 0 is the layout for the title slide
    title_slide = presentation.slides.add_slide(title_slide_layout)
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]

    title.text = data["title"]
    subtitle.text = data["subtitle"]

    # Set dark mode colors for the title slide
    set_background_color(title_slide, (23, 23, 23))  # In RGB
    set_text_color_and_size(title_slide, (236, 236, 236), 36, 24)

    # Add Other Slides
    content_slide_layout = presentation.slide_layouts[
        1
    ]  # 1 is the layout for content slides
    for slide_data in data["slides"]:
        content_slide = presentation.slides.add_slide(content_slide_layout)
        slide_title = content_slide.shapes.title
        slide_content = content_slide.placeholders[1]

        slide_title.text = slide_data["title"]
        slide_content.text = slide_data["content"]

        # Set dark mode colors for the content slide
        set_background_color(content_slide, (23, 23, 23))  # In RGB
        set_text_color_and_size(content_slide, (236, 236, 236), 36, 24)

    # Save the presentation to a file
    presentation.save(f"{filename}.pptx")
